"""
Document text extraction for PDF, DOCX, XLSX, PPTX.

Each extractor returns a list of chunk dicts:
    {"text": str, "page": int | None}

The chunking is character-based with overlap so no chunk crosses
a token-budget limit when embedding.
"""
from __future__ import annotations
from pathlib import Path
from typing import Callable, List, Dict, Any, Optional, Iterable

from app.config import CHUNK_SIZE, CHUNK_OVERLAP

# Callback invoked as on_page(done_pages, total_pages) during PDF extraction.
PageCallback = Optional[Callable[[int, int], None]]
W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
W = f"{{{W_NS}}}"


# ── Chunker ───────────────────────────────────────────────────────────────────

def _chunk(text: str) -> List[str]:
    """Split text into overlapping fixed-size character chunks."""
    text = text.strip()
    if not text:
        return []
    chunks, start = [], 0
    while start < len(text):
        end   = start + CHUNK_SIZE
        piece = text[start:end].strip()
        if piece:
            chunks.append(piece)
        if end >= len(text):
            break
        start = end - CHUNK_OVERLAP
    return chunks


# ── Per-format extractors ─────────────────────────────────────────────────────

def _extract_pdf(path: Path, on_page: PageCallback = None) -> List[Dict[str, Any]]:
    import fitz  # pymupdf
    results = []
    doc = fitz.open(str(path))
    total = doc.page_count
    try:
        for page_num, page in enumerate(doc, 1):
            text = page.get_text("text")
            for chunk in _chunk(text):
                results.append({"text": chunk, "page": page_num})
            if on_page is not None:
                try:
                    on_page(page_num, total)
                except Exception:
                    pass
    finally:
        doc.close()
    return results


def _extract_docx(path: Path) -> List[Dict[str, Any]]:
    from docx import Document
    doc = Document(str(path))

    pages: dict[int, List[str]] = {}
    page = 1
    for block in _iter_docx_text_blocks(doc):
        page, segments = _docx_segments_for_page(block, page)
        for segment_page, text in segments:
            if text.strip():
                pages.setdefault(segment_page, []).append(text.strip())
        if _docx_paragraph_starts_next_page(block):
            page += 1

    results: List[Dict[str, Any]] = []
    for page_num, parts in pages.items():
        for chunk in _chunk("\n".join(parts)):
            results.append({"text": chunk, "page": page_num})
    return results


def _iter_docx_text_blocks(doc: Any) -> Iterable[Any]:
    """Yield paragraphs in document order, including paragraphs inside tables."""
    for child in doc.element.body.iterchildren():
        if child.tag == W + "p":
            yield child
        elif child.tag == W + "tbl":
            yield from child.iter(W + "p")


def _docx_segments_for_page(element: Any, page: int) -> tuple[int, List[tuple[int, str]]]:
    """Split a DOCX paragraph into text segments keyed by Word page breaks."""
    segments: List[tuple[int, str]] = []
    buf: List[str] = []

    def flush() -> None:
        text = "".join(buf).strip()
        if text:
            segments.append((page, text))
        buf.clear()

    for node in element.iter():
        if node.tag == W + "lastRenderedPageBreak":
            flush()
            page += 1
        elif node.tag == W + "br" and node.get(W + "type") == "page":
            flush()
            page += 1
        elif node.tag == W + "tab":
            buf.append("\t")
        elif node.tag in (W + "cr", W + "br"):
            buf.append("\n")
        elif node.tag == W + "t" and node.text:
            buf.append(node.text)

    flush()
    return page, segments


def _docx_paragraph_starts_next_page(element: Any) -> bool:
    sect_pr = element.find(f"./{W}pPr/{W}sectPr")
    if sect_pr is None:
        return False
    section_type = sect_pr.find(f"./{W}type")
    value = section_type.get(W + "val") if section_type is not None else None
    return value in (None, "nextPage", "oddPage", "evenPage")


def _extract_xlsx(path: Path) -> List[Dict[str, Any]]:
    import openpyxl
    results = []
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    for sheet_name in wb.sheetnames:
        ws   = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            row_text = "\t".join(
                str(cell) if cell is not None else "" for cell in row
            )
            if row_text.strip():
                rows.append(row_text)
        sheet_text = f"[Sheet: {sheet_name}]\n" + "\n".join(rows)
        for chunk in _chunk(sheet_text):
            results.append({"text": chunk, "page": None})
    wb.close()
    return results


def _extract_pptx(path: Path) -> List[Dict[str, Any]]:
    from pptx import Presentation
    results = []
    prs = Presentation(str(path))
    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        slide_text = "\n".join(texts)
        for chunk in _chunk(slide_text):
            results.append({"text": chunk, "page": slide_num})
    return results


# ── Public API ────────────────────────────────────────────────────────────────

SUPPORTED_EXTENSIONS: set[str] = {".pdf", ".docx", ".xlsx", ".pptx"}

_EXTRACTORS = {
    ".pdf":  _extract_pdf,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".pptx": _extract_pptx,
}


def extract(path: Path, on_page: PageCallback = None) -> List[Dict[str, Any]]:
    """
    Extract text chunks from *path*.
    Returns list of {"text": str, "page": int | None}.
    Raises ValueError for unsupported file types.

    ``on_page`` is currently only honoured by the PDF extractor (its page loop
    is by far the slowest); the other formats stream too quickly to bother.
    """
    suffix = path.suffix.lower()
    extractor = _EXTRACTORS.get(suffix)
    if extractor is None:
        raise ValueError(f"Unsupported file type: {suffix!r}")
    if suffix == ".pdf":
        return extractor(path, on_page=on_page)
    return extractor(path)
